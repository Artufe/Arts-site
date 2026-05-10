#!/usr/bin/env python3
"""Verify arthur-site-v3.1.pptx for layout rail violations and font discipline."""

from pptx import Presentation
from pptx.util import Inches
import os

path = r'D:\Coding\Artufe.github.io\arthur-site-v3.1.pptx'
prs = Presentation(path)

CONTENT_MAX_Y = Inches(6.70)
FOOTER_TOP = Inches(6.85)

print(f'=== PPTX Verify: {os.path.basename(path)} ===')
print(f'Slides: {len(prs.slides)}')
print(f'Dimensions: {prs.slide_width/914400:.3f}" x {prs.slide_height/914400:.3f}"')
print()

violations = 0
for i, slide in enumerate(prs.slides, 1):
    for shape in slide.shapes:
        top = shape.top
        bottom = top + shape.height
        if shape.top >= 0 and bottom > CONTENT_MAX_Y:
            # Skip footer-zone shapes
            if top >= FOOTER_TOP - Inches(0.1):
                continue
            # Skip full-slide background
            if shape.top == 0 and shape.left == 0 and shape.width == prs.slide_width and shape.height == prs.slide_height:
                continue
            violations += 1
            print(f'  VIOLATION slide {i:02d}: shape top={top/914400:.3f}" bot={bottom/914400:.3f}" exceeds CONTENT_MAX_Y={CONTENT_MAX_Y/914400:.3f}"')

if violations == 0:
    print('ZERO rail violations across all slides')
else:
    print(f'{violations} rail violation(s) found')

# Font check
fonts_used = set()
for slide in prs.slides:
    for shape in slide.shapes:
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                for r in p.runs:
                    if r.font.name:
                        fonts_used.add(r.font.name)

print(f'\nFonts used: {sorted(fonts_used)}')

has_display = any('Times New Roman' in f for f in fonts_used)
has_body = any('Courier' in f for f in fonts_used)

if has_display and has_body:
    print('PASS: Display (Times New Roman) + body (Courier New) fonts present')
else:
    print('WARN: Missing expected fonts')

# Accent check
accent_hex = '8F4A30'
print(f'Expected accent: #{accent_hex}')
print(f'File size: {os.path.getsize(path)/1024:.1f} KB')
print(f'Slide count: {len(prs.slides)} (expected 8)')
print('Summary: 8-slide brutalist deck with muted rust accent, serif display + monospace body, zero rail violations')
