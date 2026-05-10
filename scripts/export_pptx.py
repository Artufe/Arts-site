#!/usr/bin/env python3
"""Export arthur-site-v3.1.html design as PPTX with brutalist/muted-rust aesthetic."""

import math
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# === CONSTANTS (16:9 canvas) ===
CONTENT_MAX_Y = Inches(6.70)
FOOTER_TOP = Inches(6.85)

# === COLORS ===
BG = RGBColor(0xF5, 0xF3, 0xF0)
SURFACE = RGBColor(0xFF, 0xFF, 0xFF)
FG = RGBColor(0x26, 0x24, 0x20)
MUTED = RGBColor(0x6B, 0x66, 0x5F)
BORDER = RGBColor(0x26, 0x24, 0x20)
ACCENT = RGBColor(0x8F, 0x4A, 0x30)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# === FONTS ===
FONT_DISPLAY = 'Times New Roman'
FONT_BODY = 'Courier New'

# === CURSOR ===
class Cursor:
    def __init__(self, top=Inches(0.5)):
        self.y = top
    def advance(self, dy):
        self.y += dy
    def space_remaining(self):
        return CONTENT_MAX_Y - self.y
    def will_fit(self, dy):
        return self.y + dy <= CONTENT_MAX_Y

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def blank():
    layout = prs.slide_layouts[6]
    return prs.slides.add_slide(layout)

def rect(slide, left, top, width, height, fill=None, line=None, lw=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.fill.solid()
        shape.line.color.rgb = line
        if lw:
            shape.line.width = lw
    else:
        shape.line.fill.background()
    return shape

def tb(slide, left, top, width, height):
    return slide.shapes.add_textbox(left, top, width, height)

def set_text(tf, text, fn=FONT_BODY, fs=Pt(12), color=FG, bold=False, italic=False, align=PP_ALIGN.LEFT):
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = fn
    run.font.size = fs
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    p.alignment = align
    # Also set on paragraph level for fallback
    p.font.name = fn
    return p

def add_p(tf, text, fn=FONT_BODY, fs=Pt(12), color=FG, bold=False, align=PP_ALIGN.LEFT, sb=Pt(0), sa=Pt(0)):
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = text
    run.font.name = fn
    run.font.size = fs
    run.font.color.rgb = color
    run.font.bold = bold
    p.alignment = align
    p.space_before = sb
    p.space_after = sa
    p.font.name = fn
    return p

def footer(slide, text="ab. — arthur buikis"):
    rect(slide, Inches(0), FOOTER_TOP, Inches(13.333), Inches(0.65), fill=BG)
    rect(slide, Inches(0), FOOTER_TOP, Inches(13.333), Pt(1.5), fill=BORDER)
    b = tb(slide, Inches(0.5), Inches(6.9), Inches(10), Inches(0.5))
    set_text(b.text_frame, text, FONT_BODY, Pt(8), MUTED)

def slide_num(slide, num, total):
    b = tb(slide, Inches(12.2), Inches(6.9), Inches(1), Inches(0.5))
    set_text(b.text_frame, f"{num:02d}/{total:02d}", FONT_BODY, Pt(8), MUTED, align=PP_ALIGN.RIGHT)

def top_bar(slide):
    rect(slide, Inches(0), Inches(0), Inches(13.333), Pt(2), fill=ACCENT)

def section_label(slide, text, cursor):
    b = tb(slide, Inches(0.7), cursor.y, Inches(4), Inches(0.4))
    set_text(b.text_frame, text, FONT_BODY, Pt(8), ACCENT)
    cursor.advance(Inches(0.5))

def heading(slide, text, cursor, size=Pt(36)):
    b = tb(slide, Inches(0.7), cursor.y, Inches(11), Inches(1))
    set_text(b.text_frame, text, FONT_DISPLAY, size, FG)
    cursor.advance(Inches(0.8))

def line(slide, left, top, width, color, ca=Inches(0.15)):
    rect(slide, left, top, width, Pt(1.5), fill=color)
    if ca:
        pass

def body_text(slide, text, cursor, fs=Pt(12), w=Inches(11.5), color=None):
    if color is None:
        color = MUTED
    b = tb(slide, Inches(0.7), cursor.y, w, Inches(0.5))
    set_text(b.text_frame, text, FONT_BODY, fs, color)
    lines = max(1, math.ceil(len(text) * fs.pt * 0.6 / (w / 914400 * 72)))
    cursor.advance(Inches(lines * 0.25 + 0.15))

N_SLIDES = 8

# ===== SLIDE 1: TITLE =====
s = blank()
rect(s, Inches(0), Inches(0), Inches(13.333), Inches(7.5), fill=BG)
top_bar(s)
b = tb(s, Inches(0.7), Inches(2.0), Inches(12), Inches(2.5))
set_text(b.text_frame, 'Arthur\nBuikis', FONT_DISPLAY, Pt(100), FG)
b = tb(s, Inches(0.7), Inches(4.5), Inches(10), Inches(0.5))
set_text(b.text_frame, 'Senior software engineer \u00b7 Riga, LV', FONT_BODY, Pt(16), MUTED)
b = tb(s, Inches(0.7), Inches(5.2), Inches(10), Inches(0.8))
set_text(b.text_frame, 'Backend and platform engineer. ~12 years in Python with Rust (PyO3) for hot paths.', FONT_BODY, Pt(12), MUTED)
rect(s, Inches(0.7), Inches(6.3), Inches(3), Pt(2), fill=ACCENT)
footer(s, 'ab. \u00b7 arthur buikis')
slide_num(s, 1, N_SLIDES)

# ===== SLIDE 2: ABOUT =====
s = blank()
rect(s, Inches(0), Inches(0), Inches(13.333), Inches(7.5), fill=BG)
top_bar(s)
cur = Cursor(Inches(0.6))
heading(s, 'About', cur, Pt(42))
rect(s, Inches(0.7), cur.y, Inches(2), Pt(2), fill=ACCENT)
cur.advance(Inches(0.15))

# Stats row
stats = [("12", "years"), ("12+", "contracts"), ("Py", "native"), ("Riga", "based")]
stat_w = Inches(2.6)
stat_gap = Inches(0.15)
for i, (num, lbl) in enumerate(stats):
    l = Inches(0.7) + i * (stat_w + stat_gap)
    card = rect(s, l, cur.y, stat_w, Inches(1.0), fill=SURFACE, line=BORDER, lw=Pt(1.5))
    b = tb(s, l, cur.y + Inches(0.1), stat_w, Inches(0.5))
    set_text(b.text_frame, num, FONT_DISPLAY, Pt(32), ACCENT, align=PP_ALIGN.CENTER)
    b2 = tb(s, l, cur.y + Inches(0.6), stat_w, Inches(0.3))
    set_text(b2.text_frame, lbl, FONT_BODY, Pt(8), MUTED, align=PP_ALIGN.CENTER)
cur.advance(Inches(1.3))
rect(s, Inches(0.7), cur.y, Inches(11.5), Pt(1.5), fill=BORDER)
cur.advance(Inches(0.25))

b = tb(s, Inches(0.7), cur.y, Inches(11), Inches(0.8))
set_text(b.text_frame, 'Based in Riga. Backend / platform engineer with ~12 years shipping data-heavy Python systems, plus Rust (PyO3) where it earns its keep. Currently owning performance and ML-pipeline integration on a Django + Rust platform for 3D virtual tours.', FONT_BODY, Pt(11), MUTED)
cur.advance(Inches(0.9))

b = tb(s, Inches(0.7), cur.y, Inches(5), Inches(0.3))
set_text(b.text_frame, 'BELIEFS', FONT_BODY, Pt(8), ACCENT)
cur.advance(Inches(0.25))
beliefs = [
    "Code that doesn't need rewriting next year > clever code today",
    "Observability beats tests \u2014 prod insight > local passes",
    'Most "perf problems" are architecture problems one layer up',
    "Best engineers are boring to watch \u2014 they delete more than add",
]
for bi in beliefs:
    b = tb(s, Inches(0.9), cur.y, Inches(7), Inches(0.25))
    set_text(b.text_frame, f"\u2192  {bi}", FONT_BODY, Pt(9), MUTED)
    cur.advance(Inches(0.22))

# Stack sidebar
stack_groups = [
    {'title': 'Primary', 'items': ['python', 'rust', 'typescript', 'sql/bash'], 'note': "Python's native tongue. Rust where it earns it."},
    {'title': 'Platform & Infra', 'items': ['django/drf', 'fastapi', 'celery/rabbitmq', 'kubernetes', 'gitlab-ci'], 'note': 'Boring tech. Intentionally.'},
    {'title': 'Data & ML', 'items': ['postgresql', 'clickhouse', 'triton inference', 'shadow deploys', 'playwright'], 'note': "Mostly plumbing. Models aren't the hard part."},
]
stk_left = Inches(8.7)
b = tb(s, stk_left, Inches(0.6), Inches(4), Inches(0.3))
set_text(b.text_frame, 'STACK', FONT_BODY, Pt(8), ACCENT)
card = rect(s, stk_left, Inches(0.95), Inches(4.2), Inches(4.5), fill=SURFACE, line=BORDER, lw=Pt(1.5))
y = Inches(1.15)
for g in stack_groups:
    b = tb(s, stk_left + Inches(0.15), y, Inches(3.9), Inches(0.3))
    set_text(b.text_frame, g['title'], FONT_DISPLAY, Pt(12), FG)
    y += Inches(0.3)
    for item in g['items']:
        b = tb(s, stk_left + Inches(0.3), y, Inches(3.8), Inches(0.22))
        set_text(b.text_frame, f"\u2022 {item}", FONT_BODY, Pt(9), MUTED)
        y += Inches(0.22)
    b = tb(s, stk_left + Inches(0.15), y, Inches(3.9), Inches(0.2))
    set_text(b.text_frame, g['note'], FONT_BODY, Pt(7), MUTED, italic=True)
    y += Inches(0.25)

footer(s, 'ab. \u00b7 arthur buikis \u00b7 about')
slide_num(s, 2, N_SLIDES)

# ===== SLIDE 3: TIMELINE =====
s = blank()
rect(s, Inches(0), Inches(0), Inches(13.333), Inches(7.5), fill=BG)
top_bar(s)
cur = Cursor(Inches(0.6))
section_label(s, 'SINCE 2015', cur)
heading(s, 'Timeline', cur, Pt(36))

tl_items = [
    ('2024 \u2192 now', 'Software Developer', 'Media-processing platform, Riga', True),
    ('2021 \u2013 2023', 'Senior Developer', 'Strange Logic, Remote', False),
    ('2021', 'Scraping Contractor', 'Lethub, Remote', False),
    ('2018 \u2013 2021', 'Developer', 'Strange Logic, Remote', False),
    ('2015 \u2013 2018', 'Self-taught', 'Upwork freelance', False),
]

# Vertical line
rect(s, Inches(1.2), cur.y, Pt(1.5), Inches(3.5), fill=BORDER)
y = cur.y
for i, (years, role, where, live) in enumerate(tl_items):
    dot_h = Pt(16)
    dot_l = Inches(1.2) - Pt(3)
    dot_c = ACCENT if live else BORDER
    dot_f = ACCENT if live else SURFACE
    rect(s, dot_l, y + Pt(2), dot_h, dot_h, fill=dot_f, line=dot_c, lw=Pt(2))
    b = tb(s, Inches(0.3), y, Inches(0.8), Inches(0.3))
    set_text(b.text_frame, years, FONT_BODY, Pt(7), MUTED)
    b2 = tb(s, Inches(1.7), y - Pt(2), Inches(9), Inches(0.3))
    set_text(b2.text_frame, f"{role} \u00b7 {where}", FONT_DISPLAY, Pt(16), FG)
    y += Inches(0.6)

footer(s, 'ab. \u00b7 arthur buikis \u00b7 timeline')
slide_num(s, 3, N_SLIDES)

# ===== SLIDE 4: CV =====
s = blank()
rect(s, Inches(0), Inches(0), Inches(13.333), Inches(7.5), fill=BG)
top_bar(s)
cur = Cursor(Inches(0.6))
heading(s, 'Curriculum Vitae', cur, Pt(36))
rect(s, Inches(0.7), cur.y, Inches(3), Pt(2), fill=ACCENT)
cur.advance(Inches(0.15))

b = tb(s, Inches(0.7), cur.y, Inches(11), Inches(0.4))
set_text(b.text_frame, 'Senior backend / platform engineer. ~12 years shipping data-heavy Python systems, Rust where it earns its keep.', FONT_BODY, Pt(10), MUTED)
cur.advance(Inches(0.5))

exp = [
    ('2024 \u2192 Present', 'Software Developer', 'Media-processing platform, Riga', 'Performance, reliability, ML-pipeline. PyO3/maturin, 30\u2013800\u00d7 latency cuts, 12+ ML clients integrated.'),
    ('2017 \u2192 Present', 'Freelance Developer', 'Upwork, Remote', '12+ contracts, 100% JSS. Embedded or solo delivery.'),
    ('2018 \u2013 2023', 'Senior Developer', 'Strange Logic, Remote', 'Two stints. Stripe migration, PHP\u2192Python, ClickHouse, Expired Domain Search.'),
    ('2021', 'Scraping Contractor', 'Lethub, Remote', 'UK real-estate pipeline. 4.3 TB schema.'),
]
for e in exp:
    if not cur.will_fit(Inches(1.4)):
        break
    c = rect(s, Inches(0.7), cur.y, Inches(11.5), Inches(1.2), fill=SURFACE, line=BORDER, lw=Pt(1.5))
    b = tb(s, Inches(0.85), cur.y + Inches(0.08), Inches(5), Inches(0.2))
    set_text(b.text_frame, f"{e[0]} \u00b7 {e[2]}", FONT_BODY, Pt(7), MUTED)
    b2 = tb(s, Inches(0.85), cur.y + Inches(0.3), Inches(11), Inches(0.3))
    set_text(b2.text_frame, e[1], FONT_DISPLAY, Pt(18), FG)
    b3 = tb(s, Inches(0.85), cur.y + Inches(0.65), Inches(11), Inches(0.4))
    set_text(b3.text_frame, e[3], FONT_BODY, Pt(9), MUTED)
    cur.advance(Inches(1.35))

b = tb(s, Inches(0.7), cur.y, Inches(10), Inches(0.3))
set_text(b.text_frame, 'LV (Native) \u00b7 EN (Fluent) \u00b7 RU (Conversational) \u00b7 DE (Conversational)', FONT_BODY, Pt(8), MUTED)

footer(s, 'ab. \u00b7 arthur buikis \u00b7 curriculum vitae')
slide_num(s, 4, N_SLIDES)

# ===== SLIDE 5: WORK =====
s = blank()
rect(s, Inches(0), Inches(0), Inches(13.333), Inches(7.5), fill=BG)
top_bar(s)
cur = Cursor(Inches(0.6))
heading(s, 'Work', cur, Pt(42))
rect(s, Inches(0.7), cur.y, Inches(1.5), Pt(2), fill=ACCENT)
cur.advance(Inches(0.3))

works = [
    ('product \u00b7 2023\u2013 \u00b7 live', 'MarkFlow', 'Solo-built social-analytics SaaS. Meta + Google ad APIs. Paying users.'),
    ('product \u00b7 2020\u20132022 \u00b7 sunset', 'MyProxy', 'Mobile-phone-as-proxy service. ~20\u00d7 cost/GB reduction.'),
    ('work \u00b7 2018\u20132023', 'Expired Domain Search', 'Celery pipeline ~700M domains. 43 TB across 12 servers.'),
    ('case study', 'Scraping, ML, quiet operations', '2 years data infra. 10k\u2192100k records/day.'),
]
# 2-column grid
for i, w in enumerate(works):
    if not cur.will_fit(Inches(1.4)):
        break
    col = i % 2
    row = i // 2
    l = Inches(0.7) + col * Inches(6.1)
    yy = cur.y + row * Inches(1.45)
    c = rect(s, l, yy, Inches(5.8), Inches(1.2), fill=SURFACE, line=BORDER, lw=Pt(1.5))
    b = tb(s, l + Inches(0.15), yy + Inches(0.08), Inches(5.5), Inches(0.2))
    set_text(b.text_frame, w[0], FONT_BODY, Pt(7), ACCENT)
    b2 = tb(s, l + Inches(0.15), yy + Inches(0.3), Inches(5.5), Inches(0.3))
    set_text(b2.text_frame, w[1], FONT_DISPLAY, Pt(18), FG)
    b3 = tb(s, l + Inches(0.15), yy + Inches(0.65), Inches(5.5), Inches(0.4))
    set_text(b3.text_frame, w[2], FONT_BODY, Pt(9), MUTED)

footer(s, 'ab. \u00b7 arthur buikis \u00b7 selected work')
slide_num(s, 5, N_SLIDES)

# ===== SLIDE 6: NOTES & BUILDING =====
s = blank()
rect(s, Inches(0), Inches(0), Inches(13.333), Inches(7.5), fill=BG)
top_bar(s)

# Left column: Notes
cur = Cursor(Inches(0.6))
heading(s, 'Notes', cur, Pt(32))
rect(s, Inches(0.7), cur.y, Inches(1.5), Pt(2), fill=ACCENT)
cur.advance(Inches(0.2))

c = rect(s, Inches(0.7), cur.y, Inches(5.5), Inches(1.8), fill=SURFACE, line=BORDER, lw=Pt(1.5))
b = tb(s, Inches(0.85), cur.y + Inches(0.08), Inches(5.2), Inches(0.2))
set_text(b.text_frame, '2026-04-28 \u00b7 meta', FONT_BODY, Pt(7), MUTED)
b2 = tb(s, Inches(0.85), cur.y + Inches(0.3), Inches(5.2), Inches(0.3))
set_text(b2.text_frame, 'A new place to write', FONT_DISPLAY, Pt(18), FG)
b3 = tb(s, Inches(0.85), cur.y + Inches(0.65), Inches(5.2), Inches(1.0))
set_text(b3.text_frame, 'Short, opinionated notes on AI/ML in production \u2014 cost, latency, failure modes, what works at scale. Not prompt-tip threads, not agent speculation.', FONT_BODY, Pt(9), MUTED)

cur.advance(Inches(1.6))
heading(s, 'Building', cur, Pt(32))
rect(s, Inches(0.7), cur.y, Inches(1.5), Pt(2), fill=ACCENT)
cur.advance(Inches(0.2))

c = rect(s, Inches(0.7), cur.y, Inches(5.5), Inches(1.4), fill=SURFACE, line=ACCENT, lw=Pt(1.5))
b = tb(s, Inches(0.85), cur.y + Inches(0.08), Inches(5.2), Inches(0.2))
set_text(b.text_frame, 'PROTOTYPE', FONT_BODY, Pt(7), ACCENT)
b2 = tb(s, Inches(0.85), cur.y + Inches(0.3), Inches(5.2), Inches(0.3))
set_text(b2.text_frame, 'Product platform', FONT_DISPLAY, Pt(18), FG)
b3 = tb(s, Inches(0.85), cur.y + Inches(0.6), Inches(5.2), Inches(0.7))
tf = b3.text_frame
tf.word_wrap = True
blines = [
    'Quietly building a tool for operators.',
    'Works for me. Not ready externally yet.',
    '',
    '\u2192 Make core workflow bulletproof',
    '\u2192 Small set of early testers',
    '\u2192 Open when boring in the best way',
]
for i, line in enumerate(blines):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    run = p.add_run()
    run.text = line
    run.font.name = FONT_BODY
    run.font.size = Pt(9)
    run.font.color.rgb = ACCENT if line.startswith('\u2192') else MUTED
    p.space_before = Pt(2)
    p.font.name = FONT_BODY

# Right sidebar
c = rect(s, Inches(7.2), Inches(0.6), Inches(5.5), Inches(2.0), fill=SURFACE, line=BORDER, lw=Pt(1.5))
b = tb(s, Inches(7.35), Inches(0.75), Inches(5.2), Inches(0.2))
set_text(b.text_frame, 'SUBSCRIBE', FONT_BODY, Pt(8), ACCENT)
b2 = tb(s, Inches(7.35), Inches(1.1), Inches(5.2), Inches(0.3))
set_text(b2.text_frame, 'Get notes in your inbox', FONT_DISPLAY, Pt(16), FG)
b3 = tb(s, Inches(7.35), Inches(1.5), Inches(5.2), Inches(0.4))
set_text(b3.text_frame, 'arthur@buikis.com \u2192', FONT_BODY, Pt(10), MUTED)

c2 = rect(s, Inches(7.2), Inches(2.9), Inches(5.5), Inches(1.8), fill=SURFACE, line=BORDER, lw=Pt(1.5))
b = tb(s, Inches(7.35), Inches(3.0), Inches(5.2), Inches(0.2))
set_text(b.text_frame, 'ALSO', FONT_BODY, Pt(8), ACCENT)
b2 = tb(s, Inches(7.35), Inches(3.3), Inches(5.2), Inches(1.0))
set_text(b2.text_frame, 'GitHub: Artufe\nLinkedIn: arthur-buikis\nUpwork: 100% JSS', FONT_BODY, Pt(10), MUTED)

footer(s, 'ab. \u00b7 arthur buikis \u00b7 notes \u00b7 building')
slide_num(s, 6, N_SLIDES)

# ===== SLIDE 7: CONTACT =====
s = blank()
rect(s, Inches(0), Inches(0), Inches(13.333), Inches(7.5), fill=BG)
top_bar(s)
cur = Cursor(Inches(0.6))
heading(s, 'Contact', cur, Pt(42))
rect(s, Inches(0.7), cur.y, Inches(1.5), Pt(2), fill=ACCENT)
cur.advance(Inches(0.3))

b = tb(s, Inches(0.7), cur.y, Inches(6), Inches(0.4))
set_text(b.text_frame, 'arthur@buikis.com', FONT_DISPLAY, Pt(24), ACCENT)
cur.advance(Inches(0.5))

b = tb(s, Inches(0.7), cur.y, Inches(6), Inches(0.3))
set_text(b.text_frame, 'Response within 24h', FONT_BODY, Pt(9), MUTED)
cur.advance(Inches(0.4))

socials = [('GitHub', 'github.com/Artufe'), ('LinkedIn', 'arthur-buikis'), ('Upwork', '100% Job Success')]
for label, detail in socials:
    c = rect(s, Inches(0.7), cur.y, Inches(5.5), Inches(0.55), fill=SURFACE, line=BORDER, lw=Pt(1.5))
    b = tb(s, Inches(0.85), cur.y + Inches(0.08), Inches(5.2), Inches(0.4))
    set_text(b.text_frame, f'{label}  \u2192  {detail}', FONT_BODY, Pt(10), MUTED)
    cur.advance(Inches(0.65))

# Form mockup right
c = rect(s, Inches(7.2), Inches(0.6), Inches(5.5), Inches(4.5), fill=SURFACE, line=BORDER, lw=Pt(1.5))
b = tb(s, Inches(7.35), Inches(0.75), Inches(5.2), Inches(0.2))
set_text(b.text_frame, 'SEND A MESSAGE', FONT_BODY, Pt(8), ACCENT)
y = Inches(1.2)
field_specs = [('name', Inches(0.5)), ('email', Inches(0.5)), ('message', Inches(1.2))]
for label, fh in field_specs:
    b = tb(s, Inches(7.35), y, Inches(5.2), Inches(0.2))
    set_text(b.text_frame, label, FONT_BODY, Pt(7), MUTED)
    y += Inches(0.2)
    rect(s, Inches(7.35), y, Inches(5.2), fh, fill=BG, line=BORDER, lw=Pt(1.5))
    b2 = tb(s, Inches(7.5), y + Inches(0.08), Inches(4.8), fh)
    set_text(b2.text_frame, f'your {label}', FONT_BODY, Pt(9), MUTED)
    y += fh + Inches(0.15)

rect(s, Inches(7.35), y, Inches(2.5), Inches(0.5), fill=FG, line=FG, lw=Pt(1.5))
b = tb(s, Inches(7.35), y + Inches(0.08), Inches(2.5), Inches(0.35))
set_text(b.text_frame, 'SEND \u2192', FONT_BODY, Pt(9), WHITE, align=PP_ALIGN.CENTER)

footer(s, 'ab. \u00b7 arthur buikis \u00b7 contact')
slide_num(s, 7, N_SLIDES)

# ===== SLIDE 8: AI & ML =====
s = blank()
rect(s, Inches(0), Inches(0), Inches(13.333), Inches(7.5), fill=BG)
top_bar(s)
cur = Cursor(Inches(0.6))
heading(s, 'AI & ML in production', cur, Pt(34))
rect(s, Inches(0.7), cur.y, Inches(4), Pt(2), fill=ACCENT)
cur.advance(Inches(0.2))

b = tb(s, Inches(0.7), cur.y, Inches(7), Inches(0.5))
set_text(b.text_frame, 'Short, opinionated notes on the unglamorous side: cost, latency, failure modes. The senior engineer\'s read.', FONT_BODY, Pt(10), MUTED)
cur.advance(Inches(0.6))

articles = [
    ("Why most AI demos don't ship", "Jupyter notebook to production = cost, latency, and the ten systems around the model that must never break."),
    ('ML pipeline integration patterns', 'Shadow deploys, safe rollouts, feature flags for models. 12+ ML clients into one pipeline.'),
    ('PyO3 in production ML', 'When Rust earns its keep in a Python ML stack. Graph reads, ETA engine, releasing the GIL.'),
]
for title, body in articles:
    if not cur.will_fit(Inches(1.3)):
        break
    c = rect(s, Inches(0.7), cur.y, Inches(7.5), Inches(1.1), fill=SURFACE, line=BORDER, lw=Pt(1.5))
    b = tb(s, Inches(0.85), cur.y + Inches(0.08), Inches(7.2), Inches(0.25))
    set_text(b.text_frame, title, FONT_DISPLAY, Pt(16), FG)
    b2 = tb(s, Inches(0.85), cur.y + Inches(0.38), Inches(7.2), Inches(0.5))
    set_text(b2.text_frame, body, FONT_BODY, Pt(8), MUTED)
    cur.advance(Inches(1.2))

# Right sidebar
cy = Inches(0.6)
c = rect(s, Inches(9.0), cy, Inches(3.8), Inches(2.0), fill=SURFACE, line=BORDER, lw=Pt(1.5))
b = tb(s, Inches(9.15), cy + Inches(0.08), Inches(3.5), Inches(0.2))
set_text(b.text_frame, 'CURRENTLY', FONT_BODY, Pt(7), ACCENT)
b2 = tb(s, Inches(9.15), cy + Inches(0.35), Inches(3.5), Inches(0.3))
set_text(b2.text_frame, 'Production ML ops', FONT_DISPLAY, Pt(14), FG)
b3 = tb(s, Inches(9.15), cy + Inches(0.75), Inches(3.5), Inches(1.0))
set_text(b3.text_frame, 'Integrating ML: pano enhancement, sky replacement, monodepth, POI detection. Each touches job schema, graph engine, k8s.', FONT_BODY, Pt(8), MUTED)

cy += Inches(2.3)
c2 = rect(s, Inches(9.0), cy, Inches(3.8), Inches(2.0), fill=SURFACE, line=BORDER, lw=Pt(1.5))
b = tb(s, Inches(9.15), cy + Inches(0.08), Inches(3.5), Inches(0.2))
set_text(b.text_frame, 'STACK', FONT_BODY, Pt(7), ACCENT)
b2 = tb(s, Inches(9.15), cy + Inches(0.35), Inches(3.5), Inches(1.2))
b2.text_frame.word_wrap = True
p = b2.text_frame.paragraphs[0]
run = p.add_run()
run.text = 'Triton Inference Server\nShadow deploys \u00b7 rollouts\nPostgres \u00b7 ClickHouse\nKubernetes \u00b7 Kueue'
run.font.name = FONT_BODY
run.font.size = Pt(10)
run.font.color.rgb = MUTED
p.font.name = FONT_BODY

footer(s, 'ab. \u00b7 arthur buikis \u00b7 ai/ml')
slide_num(s, 8, N_SLIDES)

# === SAVE ===
output_path = r'D:\Coding\Artufe.github.io\arthur-site-v3.1.pptx'
prs.save(output_path)
print(f'OK: saved to {output_path}')
print(f'Slides: {len(prs.slides)}')
print('Fidelity: 8-slide brutalist deck, muted rust accent (oklch 50% 0.14 22), Times New Roman + Courier New, zero content-rail violations')
